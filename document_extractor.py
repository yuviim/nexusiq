import base64
import io
import logging
import zipfile
from dataclasses import dataclass, field

import httpx
import pdfplumber
import pytesseract
from PIL import Image
from docx import Document
from pptx import Presentation

logger = logging.getLogger(__name__)

_OCR_THRESHOLD = 50


@dataclass
class ExtractedDocument:
    text:     str
    pages:    int
    tables:   list = field(default_factory=list)
    warnings: list = field(default_factory=list)
    method:   str  = "native"


class DocumentExtractor:

    def extract_from_url(
        self,
        url: str,
        content_type: str = "",
        filename: str = "",
    ) -> ExtractedDocument:
        logger.info(f"Downloading: {filename or url[:80]}")
        try:
            response = httpx.get(url, follow_redirects=True, timeout=60.0)
            response.raise_for_status()
            data = response.content
        except Exception as exc:
            return ExtractedDocument(
                text="", pages=0,
                warnings=[f"Download failed: {exc}"],
                method="failed",
            )
        return self._route(data, content_type, filename)

    def extract_from_b64(
        self,
        content_b64: str,
        content_type: str = "",
        filename: str = "",
    ) -> ExtractedDocument:
        try:
            data = base64.b64decode(content_b64)
        except Exception as exc:
            return ExtractedDocument(
                text="", pages=0,
                warnings=[f"Base64 decode failed: {exc}"],
                method="failed",
            )
        return self._route(data, content_type, filename)

    def _route(
        self,
        data: bytes,
        content_type: str,
        filename: str,
    ) -> ExtractedDocument:
        ct = content_type.lower()
        fn = filename.lower()

        if "pdf" in ct or fn.endswith(".pdf"):
            return self._extract_pdf(data)
        elif "wordprocessingml" in ct or fn.endswith(".docx"):
            return self._extract_docx(data)
        elif "presentationml" in ct or fn.endswith(".pptx"):
            return self._extract_pptx(data)
        else:
            return ExtractedDocument(
                text="", pages=0,
                warnings=[f"Unsupported type: {content_type}"],
                method="unsupported",
            )

    def _extract_pdf(self, data: bytes) -> ExtractedDocument:
        pages_text = []
        tables_text = []
        warnings = []
        ocr_pages = 0
        total_pages = 0

        try:
            with pdfplumber.open(io.BytesIO(data)) as pdf:
                total_pages = len(pdf.pages)
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text() or ""

                    for table in page.extract_tables():
                        if table:
                            rows = [
                                " | ".join(str(c) if c else "" for c in row)
                                for row in table
                            ]
                            tables_text.append("\n".join(rows))

                    if len(text.strip()) < _OCR_THRESHOLD:
                        try:
                            img = page.to_image(resolution=200).original
                            text = self._ocr_image(img)
                            ocr_pages += 1
                        except Exception as exc:
                            warnings.append(f"OCR failed on page {i+1}: {exc}")

                    if text.strip():
                        pages_text.append(f"[Page {i+1}]\n{text.strip()}")

        except Exception as exc:
            return ExtractedDocument(
                text="", pages=0,
                warnings=[f"PDF extraction failed: {exc}"],
                method="failed",
            )

        method = (
            "ocr"    if ocr_pages == total_pages else
            "hybrid" if ocr_pages > 0            else
            "native"
        )
        full_text = "\n\n".join(pages_text)
        if tables_text:
            full_text += "\n\nTABLES:\n" + "\n\n".join(tables_text)

        logger.info(f"PDF: {total_pages} pages, {ocr_pages} OCR, method={method}")
        return ExtractedDocument(
            text=full_text, pages=total_pages,
            tables=tables_text, warnings=warnings, method=method,
        )

    def _extract_docx(self, data: bytes) -> ExtractedDocument:
        paragraphs = []
        tables_text = []
        warnings = []

        try:
            doc = Document(io.BytesIO(data))

            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text.strip())

            for table in doc.tables:
                rows = [
                    " | ".join(cell.text.strip() for cell in row.cells)
                    for row in table.rows
                ]
                if rows:
                    tables_text.append("\n".join(rows))

            embedded_text = self._extract_docx_embedded_pdfs(data, warnings)

        except Exception as exc:
            return ExtractedDocument(
                text="", pages=0,
                warnings=[f"DOCX extraction failed: {exc}"],
                method="failed",
            )

        parts = []
        if paragraphs:
            parts.append("\n".join(paragraphs))
        if tables_text:
            parts.append("TABLES:\n" + "\n\n".join(tables_text))
        if embedded_text:
            parts.append("EMBEDDED DOCUMENTS:\n" + "\n\n".join(embedded_text))

        logger.info(f"DOCX: {len(paragraphs)} paragraphs, {len(tables_text)} tables, {len(embedded_text)} embedded")
        return ExtractedDocument(
            text="\n\n".join(parts),
            pages=len(paragraphs),
            tables=tables_text,
            warnings=warnings,
            method="native",
        )

    def _extract_docx_embedded_pdfs(
        self,
        data: bytes,
        warnings: list,
    ) -> list:
        extracted = []
        try:
            with zipfile.ZipFile(io.BytesIO(data)) as zf:
                pdf_parts = [
                    name for name in zf.namelist()
                    if name.startswith("word/embeddings/")
                    and name.lower().endswith(".pdf")
                ]
                for path in pdf_parts:
                    try:
                        pdf_bytes = zf.read(path)
                        result = self._extract_pdf(pdf_bytes)
                        if result.text.strip():
                            extracted.append(
                                f"[Embedded: {path.split('/')[-1]}]\n{result.text}"
                            )
                            logger.info(f"Extracted embedded PDF: {path}")
                    except Exception as exc:
                        warnings.append(f"Embedded {path} failed: {exc}")
        except Exception as exc:
            warnings.append(f"ZIP read failed: {exc}")
        return extracted

    def _extract_pptx(self, data: bytes) -> ExtractedDocument:
        slides_text = []
        warnings = []

        try:
            prs = Presentation(io.BytesIO(data))
            for i, slide in enumerate(prs.slides):
                parts = []

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                parts.append(para.text.strip())
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            parts.append(" | ".join(cells))

                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        parts.append(f"[Notes] {notes}")

                if parts:
                    slides_text.append(f"[Slide {i+1}]\n" + "\n".join(parts))

        except Exception as exc:
            return ExtractedDocument(
                text="", pages=0,
                warnings=[f"PPTX extraction failed: {exc}"],
                method="failed",
            )

        logger.info(f"PPTX: {len(slides_text)} slides extracted")
        return ExtractedDocument(
            text="\n\n".join(slides_text),
            pages=len(slides_text),
            warnings=warnings,
            method="native",
        )

    def _ocr_image(self, image) -> str:
        if not isinstance(image, Image.Image):
            image = Image.fromarray(image)
        return pytesseract.image_to_string(image, lang="eng").strip()